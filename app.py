
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import requests
import json
EVENT_CODE = "lytrq"
FTC_API_URL = "http://localhost"
ppt = Presentation('template.pptx')
fullEvent = requests.get(f"{FTC_API_URL}/api/v2/events/{EVENT_CODE}/full/").json()
# fullEvent = json.load(open('fullEvent.json'))

def getSeriesName(series):
    if series == 1:
        return "1st"
    elif series == 2:
        return "2nd"
    else:
        return "3rd"
def CreateNewSlideWithTextCenter(ppt,Maintext,SubText = None):
    slide = ppt.slides.add_slide(ppt.slide_layouts[15])
    placeholder = slide.placeholders[0]  
    text_frame = placeholder.text_frame
    p = text_frame.add_paragraph()
    p.text = Maintext
    p.alignment = PP_ALIGN.CENTER
    if SubText is not None:
            placeholder = slide.placeholders[1]  
            text_frame = placeholder.text_frame
            p = text_frame.add_paragraph()
            p.text = SubText
            p.alignment = PP_ALIGN.CENTER

new_teams_dict = {}
for team in fullEvent["teamList"]["teams"]:
    team_number = team.pop("number")  # This removes and returns the number
    new_teams_dict[str(team_number)] = team
fullEvent["teamList"]["teams"] = new_teams_dict

EventTitle = fullEvent["event"]["name"]
CreateNewSlideWithTextCenter(ppt,EventTitle)

# Add a slide for each team

for award in fullEvent["awardList"]["awards"]:
    if award["isTeamAward"] == True and award["winners"] != []: 
        CreateNewSlideWithTextCenter(ppt,award["name"])
        for team in reversed(award["winners"]):
            if(team['team'] == -1 or team['team'] == 0):
                continue
            CreateNewSlideWithTextCenter(ppt,f"{getSeriesName(team['series'])} {award['name']}",f"Team #{team['team']} \n {fullEvent['teamList']['teams'][str(team['team'])]['name']}")

ppt.save(f'{EventTitle} Awards.pptx')
